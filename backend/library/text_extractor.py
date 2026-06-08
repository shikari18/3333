import io
import logging
from typing import Dict, Any

logger = logging.getLogger('nitemind')

def extract_text_from_bytes(file_bytes: bytes, extension: str) -> Dict[str, Any]:
    """
    Extract text from common document formats.
    Returns: {'text': str, 'status': 'success'|'error', 'error': str}
    """
    content = {'text': '', 'status': 'success'}
    
    try:
        ext = extension.lower()
        if ext == '.pdf':
            from .pdf_extractor import extract_pdf_content
            pdf_data = extract_pdf_content(file_bytes=file_bytes)
            content['text'] = pdf_data['text']
            content['pdf_data'] = pdf_data # Keep images etc for caller
        
        elif ext in ['.docx', '.doc']:
            import docx
            doc = docx.Document(io.BytesIO(file_bytes))
            full_text = []
            for para in doc.paragraphs:
                full_text.append(para.text)
            content['text'] = '\n'.join(full_text)
            
        elif ext in ['.txt', '.md', '.py', '.js', '.ts', '.css', '.html']:
            content['text'] = file_bytes.decode('utf-8', errors='ignore')

        elif ext in ['.pptx', '.ppt']:
            try:
                from pptx import Presentation
                from pptx.util import Inches
                import io as _io

                prs = Presentation(_io.BytesIO(file_bytes))
                slides_text = []
                extracted_images = []

                for i, slide in enumerate(prs.slides):
                    slide_parts = [f"--- Slide {i+1} ---"]

                    # Extract text from all shapes
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            slide_parts.append(shape.text.strip())
                        # Extract images from slides
                        if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                            try:
                                img_blob = shape.image.blob
                                img_ext = shape.image.ext
                                extracted_images.append({
                                    'data': img_blob,
                                    'page': i + 1,
                                    'ext': img_ext or 'png',
                                    'is_large': len(img_blob) > 50000,
                                })
                            except Exception:
                                pass

                    # Extract speaker notes
                    try:
                        notes_slide = slide.notes_slide
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_parts.append(f"[Speaker Notes: {notes_text}]")
                    except Exception:
                        pass

                    if len(slide_parts) > 1:
                        slides_text.append('\n'.join(slide_parts))

                content['text'] = '\n\n'.join(slides_text)
                if extracted_images:
                    content['slide_images'] = extracted_images
                    content['page_count'] = len(prs.slides)

            except ImportError:
                content['status'] = 'error'
                content['error'] = 'python-pptx not installed'
            except Exception as e:
                content['status'] = 'error'
                content['error'] = str(e)
            
        else:
            content['status'] = 'error'
            content['error'] = f'Unsupported extension: {ext}'
            
    except Exception as e:
        logger.error(f'Extraction error for {extension}: {e}')
        content['status'] = 'error'
        content['error'] = str(e)
        
    return content
